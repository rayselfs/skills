#!/usr/bin/env python3
"""
Weekly Report PPTX Generator

Generates or modifies a PowerPoint weekly report from commit data.

Usage:
  # Create from scratch
  python pptx-generator.py --input commits.json --summary summary.json --output report.pptx

  # Modify existing template (uses {{PLACEHOLDER}} markers in slides)
  python pptx-generator.py --template template.pptx --input commits.json --summary summary.json --output report.pptx

  # Specify week definition (default: mon-fri)
  python pptx-generator.py ... --week-def mon-sun

Template placeholders (add these text strings inside your .pptx):
  {{SUMMARY}}                    Executive summary paragraph
  {{WEEK}}                       Week number (e.g. 23)
  {{YEAR}}                       Year (e.g. 2025)
  {{STATS}}                      "N commits across N repos"
  {{DATE_RANGE}}                 "Jun 02 – Jun 06, 2025"
  {{CATEGORY_Infrastructure}}    Narrative for category named "Infrastructure"
  {{CATEGORY_Features}}          Narrative for category named "Features"
  (any category name works — {{CATEGORY_{Name}}})

NOTE on split-run placeholders:
  PowerPoint sometimes splits a single placeholder like {{SUMMARY}} across
  multiple runs (e.g. {{SUM and MARY}}). This script merges all runs in a
  paragraph before checking for placeholders, so it handles this correctly.
"""

import argparse
import json
import sys
from datetime import date, timedelta
from pathlib import Path


def get_week_range(week_def: str = "mon-fri"):
    today = date.today()
    mon = today - timedelta(days=today.weekday())

    if week_def == "mon-fri":
        start, end = mon, mon + timedelta(days=4)
    elif week_def == "mon-sun":
        start, end = mon, mon + timedelta(days=6)
    elif week_def == "last7":
        end, start = today, today - timedelta(days=7)
    else:
        print(f"ERROR: unknown week_def '{week_def}'. Use mon-fri, mon-sun, or last7", file=sys.stderr)
        sys.exit(1)

    week_num = mon.isocalendar()[1]
    year = mon.year
    return start, end, week_num, year


def build_replacements(summary_data: dict, stats: dict, week_def: str) -> dict:
    """Build the {{PLACEHOLDER}} → value mapping."""
    start, end, week_num, year = get_week_range(week_def)
    replacements = {
        "SUMMARY": summary_data.get("executive_summary", ""),
        "WEEK": str(week_num),
        "YEAR": str(year),
        "STATS": f"{stats['total_commits']} commits across {stats['repos_touched']} repos",
        "DATE_RANGE": f"{start.strftime('%b %d')} – {end.strftime('%b %d, %Y')}",
    }
    for category, narrative in summary_data.get("categories", {}).items():
        key = f"CATEGORY_{category.replace(' ', '_')}"
        replacements[key] = narrative
    return replacements


def replace_in_paragraph(para, replacements: dict):
    """
    Replace {{KEY}} placeholders in a paragraph, handling the case where
    PowerPoint has split a single placeholder across multiple runs.

    Strategy:
    1. Collect the full text of all runs in the paragraph.
    2. If any placeholder exists in the merged text, perform replacements.
    3. Write the result back into runs[0], clear remaining runs.
       Preserves the formatting (font, size, bold, etc.) of runs[0].
    """
    if not para.runs:
        return

    full_text = "".join(r.text for r in para.runs)
    changed = False
    for key, value in replacements.items():
        marker = f"{{{{{key}}}}}"
        if marker in full_text:
            full_text = full_text.replace(marker, value)
            changed = True

    if changed:
        para.runs[0].text = full_text
        for run in para.runs[1:]:
            run.text = ""


def apply_replacements_to_slide(slide, replacements: dict):
    """Walk all paragraphs in a slide and apply replacements."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            replace_in_paragraph(para, replacements)


def modify_template(template_path: str, output_path: str, summary_data: dict, stats: dict, week_def: str):
    """Open an existing .pptx and replace {{PLACEHOLDER}} markers throughout."""
    from pptx import Presentation

    prs = Presentation(template_path)
    replacements = build_replacements(summary_data, stats, week_def)

    for slide in prs.slides:
        apply_replacements_to_slide(slide, replacements)

    prs.save(output_path)
    print(f"✓ Template modified → {output_path}")


def create_from_scratch(output_path: str, summary_data: dict, stats: dict, week_def: str):
    """Build a clean weekly report presentation from scratch."""
    from pptx import Presentation
    from pptx.util import Inches

    start, end, week_num, year = get_week_range(week_def)
    prs = Presentation()

    # Slide dimensions: widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]   # Title Slide
    content_layout = prs.slide_layouts[1]  # Title and Content

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = f"Weekly Report — Week {week_num}, {year}"
    subtitle = slide.placeholders[1]
    subtitle.text = f"{start.strftime('%B %d')} – {end.strftime('%B %d, %Y')}"

    # --- Slide 2: Executive Summary ---
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Summary"
    tf = slide.placeholders[1].text_frame
    tf.word_wrap = True
    tf.text = summary_data.get("executive_summary", "")

    # --- Slide per category ---
    for category, narrative in summary_data.get("categories", {}).items():
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = category
        tf = slide.placeholders[1].text_frame
        tf.word_wrap = True
        tf.text = narrative

    # --- Last slide: Stats ---
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Stats"
    tf = slide.placeholders[1].text_frame
    tf.text = ""
    entries = [
        ("Total commits", str(stats["total_commits"])),
        ("Repos touched", str(stats["repos_touched"])),
        ("Week", f"W{week_num} / {year}"),
        ("Period", f"{start.strftime('%Y-%m-%d')} → {end.strftime('%Y-%m-%d')}"),
    ]
    first = True
    for label, value in entries:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = f"{label}: {value}"
        first = False

    prs.save(output_path)
    print(f"✓ Created from scratch → {output_path}")


def main():
    parser = argparse.ArgumentParser(
        description="Generate or modify a weekly report .pptx from commit data"
    )
    parser.add_argument("--input", required=True, help="Path to commits.json")
    parser.add_argument(
        "--summary",
        required=True,
        help="Path to summary.json (or inline JSON string)",
    )
    parser.add_argument("--output", required=True, help="Output .pptx file path")
    parser.add_argument(
        "--template",
        help="Optional: existing .pptx template to modify instead of creating from scratch",
    )
    parser.add_argument(
        "--week-def",
        default="mon-fri",
        choices=["mon-fri", "mon-sun", "last7"],
        help="Week definition for date range (default: mon-fri)",
    )
    args = parser.parse_args()

    # Load commits
    with open(args.input) as f:
        commits = json.load(f)

    # Load summary — accept file path or raw JSON string
    summary_path = Path(args.summary)
    if summary_path.exists():
        with open(summary_path) as f:
            summary_data = json.load(f)
    else:
        try:
            summary_data = json.loads(args.summary)
        except json.JSONDecodeError:
            print("ERROR: --summary must be a valid file path or JSON string", file=sys.stderr)
            sys.exit(1)

    # Compute stats
    repos = set()
    for c in commits:
        parts = [c.get("project", ""), c.get("repo", "")]
        repos.add("/".join(p for p in parts if p))
    stats = {
        "total_commits": len(commits),
        "repos_touched": len(repos),
    }

    week_def = args.week_def

    if args.template:
        modify_template(args.template, args.output, summary_data, stats, week_def)
    else:
        create_from_scratch(args.output, summary_data, stats, week_def)


if __name__ == "__main__":
    main()
